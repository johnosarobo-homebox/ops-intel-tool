from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK   = RGBColor(0x0f, 0x11, 0x17)
PURPLE = RGBColor(0x63, 0x66, 0xf1)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xe2, 0xe8, 0xf0)
MUTED  = RGBColor(0x64, 0x74, 0x8b)
RED    = RGBColor(0xef, 0x44, 0x44)
AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
GREEN  = RGBColor(0x22, 0xc5, 0x5e)
SLATE  = RGBColor(0x1e, 0x20, 0x28)
SLATE2 = RGBColor(0x16, 0x19, 0x20)
BLUE   = RGBColor(0x60, 0xa5, 0xfa)
PURP_L = RGBColor(0xa7, 0x8b, 0xfa)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color, border_color=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def circle(slide, x, y, size, color):
    s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def pill(slide, text, x, y, w, h, bg_c, fg_c):
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = fg_c
    r.font.bold = True; r.font.name = 'Calibri'

def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = 'Calibri'

def card(slide, x, y, w, h, border_color=None):
    c = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = SLATE
    if border_color:
        c.line.color.rgb = border_color; c.line.width = Pt(1)
    else:
        c.line.color.rgb = RGBColor(0x2a,0x2d,0x3a); c.line.width = Pt(0.5)
    return c

# ── SLIDE 1: TITLE ─────────────────────────────────────────────────────────
s1 = blank_slide(prs)
bg(s1)
box(s1, 0, 0, 4.0, 7.5, RGBColor(0x09,0x0b,0x10))
bar(s1, 0, 0, 0.06, 7.5, PURPLE)
bar(s1, 4.3, 2.05, 8.8, 0.04, RGBColor(0x2a,0x2d,0x3a))

txt(s1, 'INTERNAL TOOLING  ·  SHEFFIELD OPS  ·  APRIL 2026', 4.4, 0.55, 8.5, 0.35, 8, MUTED, bold=True)
txt(s1, 'TSEG Ops', 4.4, 1.0, 8.5, 1.3, 52, WHITE, bold=True)
txt(s1, 'Operational intelligence tool', 4.4, 2.2, 8.5, 0.5, 17, MUTED)
txt(s1, 'Automates TSEG order monitoring, SLA breach\ndetection and supplier intelligence for the\nSheffield operations team.', 4.4, 2.85, 7.8, 1.1, 11.5, LIGHT)
pill(s1, 'Built independently', 4.4,  4.15, 1.75, 0.3, RGBColor(0x1e,0x1a,0x3d), PURP_L)
pill(s1, 'Production ready',    6.3,  4.15, 1.55, 0.3, RGBColor(0x0f,0x2d,0x18), GREEN)
pill(s1, '7 slides',            8.0,  4.15, 0.85, 0.3, RGBColor(0x1a,0x1d,0x26), MUTED)
txt(s1, 'John Osarobo  ·  Operational Analyst  ·  Homebox Sheffield', 4.4, 6.9, 8.5, 0.35, 9, MUTED)

# ── SLIDE 2: THE PROBLEM ───────────────────────────────────────────────────
s2 = blank_slide(prs)
bg(s2)
box(s2, 0, 0, 13.33, 1.1, SLATE2)
bar(s2, 0, 0, 0.06, 1.1, RED)
txt(s2, 'THE PROBLEM', 0.4, 0.14, 6, 0.3, 8, RED, bold=True)
txt(s2, 'Up to half a working day lost to manual checks', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

problems = [
    ('No single view', 'Trevor, Filament and the TSEG WIP sheet had to be cross-checked manually every day with no unified dashboard or consolidated output.'),
    ('SLA blind spots', 'Breach risk was invisible until already a breach. Orders in awaiting feedback were only identified through time-consuming manual review.'),
    ('Supplier patterns hidden', 'Objections from specific suppliers were buried in a flat list with no way to see which suppliers caused the most friction.'),
    ('Gas fuel errors', 'Electricity-only properties with incorrect gas services had to be identified one order at a time — slow and error-prone.'),
    ('Time cost', 'On high-volume days the combined manual workflow consumed up to half a working day per analyst, leaving less time for actual resolution.'),
]
cols = [0.3, 4.65, 8.98]
rows = [1.3, 3.85]
positions = [(0,0),(1,0),(2,0),(0,1),(1,1)]
for idx, (title, desc) in enumerate(problems):
    ci, ri = positions[idx]
    cx, cy = cols[ci], rows[ri]
    card(s2, cx, cy, 4.1, 2.2)
    circle(s2, cx+0.18, cy+0.2, 0.14, RED)
    txt(s2, title, cx+0.42, cy+0.14, 3.5, 0.35, 10.5, WHITE, bold=True)
    txt(s2, desc,  cx+0.18, cy+0.56, 3.72, 1.5, 9, MUTED)

# ── SLIDE 3: THE SOLUTION ─────────────────────────────────────────────────
s3 = blank_slide(prs)
bg(s3)
box(s3, 0, 0, 13.33, 1.1, SLATE2)
bar(s3, 0, 0, 0.06, 1.1, PURPLE)
txt(s3, 'THE SOLUTION', 0.4, 0.14, 6, 0.3, 8, PURPLE, bold=True)
txt(s3, 'Three tools. One browser tab. Five minutes.', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

tools = [
    (PURPLE, '01', 'Gas checker', 'Automatically cross-references Homebox meter data against TSEG supplier data to identify electricity-only properties incorrectly assigned a gas service.\n\nExports a clean flagged action list for remediation. Can push results directly to a dated Google Sheets tab in one click.'),
    (RED,    '02', 'SLA monitor', 'Filters all orders in awaiting feedback and calculates days elapsed since last update.\n\nRAG-rated: red for breached (8+ days), amber for at risk (6-7 days), green for within SLA. Sorted worst first. Includes supplier breakdown showing which providers have the most breaches.'),
    (AMBER,  '03', 'WIP cross-reference', 'Connects directly to the live TSEG WIP Google Sheet, reads all active tabs — Objections, Missing Meter Information, Gas Deleted, Switch Issues, API Order Errors, ET Requests, Missing Tenant Details — and joins every row to Trevor data by TSEG ID in seconds.\n\nObjections ranked and grouped by supplier with a visual bar chart.'),
]
for i, (color, num, title, desc) in enumerate(tools):
    cx = 0.3 + i * 4.35
    card(s3, cx, 1.3, 4.1, 5.85)
    bar(s3, cx, 1.3, 4.1, 0.055, color)
    txt(s3, num,   cx+0.22, 1.48, 3.7, 0.55, 26, color, bold=True)
    txt(s3, title, cx+0.22, 2.1,  3.7, 0.42, 13, WHITE, bold=True)
    txt(s3, desc,  cx+0.22, 2.62, 3.66, 4.3, 9.5, LIGHT)

# ── SLIDE 4: IMPACT ───────────────────────────────────────────────────────
s4 = blank_slide(prs)
bg(s4)
box(s4, 0, 0, 13.33, 1.1, SLATE2)
bar(s4, 0, 0, 0.06, 1.1, GREEN)
txt(s4, 'IMPACT', 0.4, 0.14, 6, 0.3, 8, GREEN, bold=True)
txt(s4, 'From half a day to five minutes', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

metrics = [
    ('~4 hrs', 'Saved per day on\nmanual TSEG checks',   GREEN),
    ('8 days', 'SLA threshold now\nmonitored automatically', RED),
    ('723',    'Orders cross-referenced\nin a single run',   AMBER),
    ('32',     'Active objections\nranked by supplier',      PURP_L),
]
for i, (val, label, color) in enumerate(metrics):
    cx = 0.3 + i * 3.25
    mc = card(s4, cx, 1.28, 3.0, 1.38)
    txt(s4, val,   cx+0.18, 1.38, 2.6, 0.65, 26, color, bold=True)
    txt(s4, label, cx+0.18, 2.0,  2.6, 0.55, 8.5, MUTED)

bar(s4, 0.3, 2.88, 12.7, 0.04, RGBColor(0x2a,0x2d,0x3a))

box(s4, 0.3,  3.05, 6.15, 4.12, SLATE, RGBColor(0x2a,0x2d,0x3a))
box(s4, 6.88, 3.05, 6.15, 4.12, RGBColor(0x0d,0x26,0x18), RGBColor(0x16,0x65,0x14))
txt(s4, 'BEFORE', 0.55,  3.22, 5.6, 0.3, 8, RED,   bold=True)
txt(s4, 'AFTER',  7.08,  3.22, 5.6, 0.3, 8, GREEN, bold=True)

befores = ['Manual cross-check of Trevor, Filament and WIP sheet daily', 'Up to half a day on high-volume periods', 'SLA breach invisible until it was already a breach', 'Supplier objection patterns not visible at a glance', 'Gas fuel errors identified one order at a time']
afters  = ['Single browser tool — drop files, hit run', 'Full analysis in under 5 minutes', 'Live RAG status on every awaiting-feedback order', 'Supplier objection rankings updated on every run', 'Electricity-only orders flagged in bulk automatically']
for i, (b, a) in enumerate(zip(befores, afters)):
    txt(s4, '–  ' + b, 0.52, 3.62 + i*0.62, 5.7, 0.52, 9, LIGHT)
    txt(s4, '+  ' + a, 7.05, 3.62 + i*0.62, 5.7, 0.52, 9, LIGHT)

# ── SLIDE 5: SUPPLIER INTELLIGENCE ────────────────────────────────────────
s5 = blank_slide(prs)
bg(s5)
box(s5, 0, 0, 13.33, 1.1, SLATE2)
bar(s5, 0, 0, 0.06, 1.1, RED)
txt(s5, 'SUPPLIER INTELLIGENCE', 0.4, 0.14, 8, 0.3, 8, RED, bold=True)
txt(s5, 'Objections ranked by supplier — escalate with precision', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

txt(s5, 'The WIP cross-reference module groups active objections by supplier and displays them as a ranked bar chart. Rather than scanning a flat list, the ops team immediately sees which suppliers are causing the most friction — enabling targeted, evidence-based escalations and more effective supplier conversations.', 0.4, 1.22, 12.5, 0.72, 10, LIGHT)

suppliers = [('E.ON Next', 11), ('British Gas', 10), ('Utility Warehouse', 4), ('Ovo Energy', 3), ('Utilita Energy', 2), ('EDF Energy', 1), ('Scottish Power', 1)]
max_v = 11
bar_area_w = 7.5
for i, (name, count) in enumerate(suppliers):
    cy = 2.18 + i * 0.59
    bw = (count / max_v) * bar_area_w
    intensity = max(0x22, int(0xef - i * 0x16))
    bar(s5, 3.1, cy+0.09, bw, 0.33, RGBColor(intensity, 0x33, 0x33))
    txt(s5, name,  0.4, cy, 2.55, 0.42, 10, WHITE if i < 2 else LIGHT, bold=(i < 2))
    txt(s5, str(count) + (' objections' if count > 1 else ' objection'), 3.1 + bw + 0.12, cy, 2.5, 0.42, 9.5, RED if i < 2 else MUTED)

txt(s5, 'The supplier WIP heat map also shows each supplier\'s full breakdown across all blocker types — objections, missing meter info, gas deleted, switch issues, missing tenant details — giving a complete operational picture of supplier friction in one view.', 0.4, 6.45, 12.5, 0.72, 9, MUTED)

# ── SLIDE 6: TECH STACK ────────────────────────────────────────────────────
s6 = blank_slide(prs)
bg(s6)
box(s6, 0, 0, 13.33, 1.1, SLATE2)
bar(s6, 0, 0, 0.06, 1.1, BLUE)
txt(s6, 'TECHNICAL OVERVIEW', 0.4, 0.14, 8, 0.3, 8, BLUE, bold=True)
txt(s6, 'Built on proven, production-grade open source', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

stack = [
    (RGBColor(0x37,0x8a,0xdd), 'Python 3.12',       'Language', 'Used by Google, NASA and most financial institutions for data tooling and automation.'),
    (PURPLE,                   'FastAPI',             'Framework', 'High-performance web framework used internally by Netflix, Uber and Microsoft.'),
    (RGBColor(0x37,0x8a,0xdd), 'pandas',             'Data',      'Industry standard for tabular data processing. Used across finance, healthcare and energy.'),
    (GREEN,                    'Google Sheets API',  'Integration','Service account authentication — Google\'s recommended enterprise method for internal tools.'),
    (PURP_L,                   'Railway',            'Hosting',   'European cloud hosting with HTTPS by default and automated GitHub deploys.'),
    (AMBER,                    'Trevor / BigQuery',  'Data source','Homebox data source. Tool consumes CSV exports — no direct database connection required.'),
]
for i, (color, name, tag, desc) in enumerate(stack):
    col = i % 3; row = i // 3
    cx = 0.3 + col * 4.35
    cy = 1.32 + row * 2.05
    card(s6, cx, cy, 4.1, 1.85, color)
    txt(s6, name, cx+0.22, cy+0.14, 2.8, 0.36, 11, color, bold=True)
    pill(s6, tag, cx+3.05, cy+0.12, 0.82, 0.26, RGBColor(0x0f,0x11,0x17), color)
    txt(s6, desc, cx+0.22, cy+0.56, 3.66, 1.1, 9, LIGHT)

bar(s6, 0.3, 5.6, 12.7, 0.04, RGBColor(0x2a,0x2d,0x3a))
security = [
    ('No data storage',    'All processing in-memory per session. Files discarded immediately after each run. Nothing persisted server-side.'),
    ('Credentials secured','Google service account key stored as encrypted environment variable — never in source code or version control.'),
    ('Access control',     'Deployed URL can be restricted to Homebox IP ranges or password-protected via Railway. HTTPS enforced by default.'),
]
for i, (title, desc) in enumerate(security):
    cx = 0.3 + i * 4.35
    txt(s6, title, cx, 5.75, 4.0, 0.32, 9, WHITE, bold=True)
    txt(s6, desc,  cx, 6.12, 4.0, 0.95, 8.5, MUTED)

# ── SLIDE 7: NEXT STEPS ───────────────────────────────────────────────────
s7 = blank_slide(prs)
bg(s7)
box(s7, 0, 0, 13.33, 1.1, SLATE2)
bar(s7, 0, 0, 0.06, 1.1, PURPLE)
txt(s7, 'NEXT STEPS', 0.4, 0.14, 6, 0.3, 8, PURPLE, bold=True)
txt(s7, 'Production ready. Proposed for team deployment.', 0.4, 0.42, 12.5, 0.6, 21, WHITE, bold=True)

steps = [
    (PURPLE, '01', 'Deploy team-wide', 'Host on Railway so all ops analysts access the tool through a shared browser link. No installation required per user. Estimated setup time: 15 minutes.\n\nAll team members would access identical functionality through the same URL, with results pushing to shared Google Sheets automatically.'),
    (AMBER,  '02', 'Daily SLA digest', 'Automated morning summary of SLA breaches and new WIP objections delivered to Slack or email each day — without needing to open the tool or run any manual checks.\n\nWould flag only newly breached or at-risk orders to avoid noise.'),
    (GREEN,  '03', 'Expand coverage', 'Extend the same data pipeline architecture to water and council tax workflows. The column auto-detection logic already handles variable export formats from different suppliers.\n\nNo rebuild required — configuration only.'),
]
for i, (color, num, title, desc) in enumerate(steps):
    cx = 0.3 + i * 4.35
    card(s7, cx, 1.32, 4.1, 5.75)
    bar(s7, cx, 1.32, 4.1, 0.055, color)
    txt(s7, num,   cx+0.22, 1.5,  3.7, 0.6,  28, color, bold=True)
    txt(s7, title, cx+0.22, 2.18, 3.7, 0.42, 13, WHITE, bold=True)
    txt(s7, desc,  cx+0.22, 2.72, 3.66, 4.1, 9.5, LIGHT)

txt(s7, 'John Osarobo  ·  Operational Analyst  ·  Homebox Sheffield  ·  April 2026', 0.4, 7.08, 10.5, 0.32, 9, MUTED)
pill(s7, 'Built independently', 9.55, 7.04, 1.75, 0.28, RGBColor(0x1e,0x1a,0x3d), PURP_L)
pill(s7, 'Production ready',    11.42, 7.04, 1.6,  0.28, RGBColor(0x0f,0x2d,0x18), GREEN)

out = os.path.expanduser('~/Desktop/TSEG_Ops_Presentation.pptx')
prs.save(out)
print(f'\nPresentation saved to: {out}')
print('7 slides | Dark theme | Ready to present')
